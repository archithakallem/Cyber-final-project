import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt

# Initialize Presentation using default template
prs = Presentation()

def add_slide(prs, layout_idx, title_text, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Set the title
    if slide.shapes.title:
        slide.shapes.title.text = title_text
    
    # Add bullet points
    if len(slide.shapes.placeholders) > 1 and len(bullet_points) > 0:
        tf = slide.shapes.placeholders[1].text_frame
        for i, point in enumerate(bullet_points):
            if isinstance(point, tuple):
                # It's a nested bullet point (Parent, [Children])
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = point[0]
                p.level = 0
                for subpoint in point[1]:
                    sub_p = tf.add_paragraph()
                    sub_p.text = subpoint
                    sub_p.level = 1
            else:
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = point
                p.level = 0
    return slide

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "CyberScan"
subtitle.text = "Security Analysis Dashboard\n\nA Unified approach to Discovery and Threat Intelligence"

# Slide 2: Problem Statement
add_slide(prs, 1, "Problem Statement", [
    "Fragmented Tooling: Teams use disconnected tools for scanning (Nmap), analysis (VirusTotal), and reporting.",
    "High Time-to-Detect: Manually correlating open ports with threat intelligence is slow.",
    "Subjective Risk: Lack of a unified, mathematical risk scoring mechanism.",
    "Poor Tracking: No built-in way to track an asset's vulnerability history over time."
])

# Slide 3: Our Solution
add_slide(prs, 1, "Our Solution", [
    "CyberScan is a unified security platform combining proactive discovery with reactive threat intelligence.",
    "Connects automated open-port detection with live VirusTotal API queries.",
    "Calculates an actionable and structured unified risk score.",
    "Visualizes the entire security posture through an interactive dashboard.",
    "Proactively sends alerts for high-risk targets."
])

# Slide 4: Objectives
add_slide(prs, 1, "Objectives", [
    "Unify Network Discovery & Threat Intelligence in one seamless workflow.",
    "Automate Risk Quantification across four dimensions (Exposure, Threat, Context, Overall).",
    "Maintain Historical Records for security drift tracking via SQLite.",
    "Provide a Responsive Dashboard for analysts to rapidly assess systems."
])

# Slide 5: Key Features Overview
add_slide(prs, 1, "Key Features Overview", [
    ("Dual Architecture", ["Scalable FastAPI Backend", "Interactive Streamlit Frontend"]),
    ("Scanner Integration", ["Nmap open-port discovery", "VirusTotal threat intelligence"]),
    ("Smart Scoring", ["Automated risk matrix calculations"]),
    ("Data Persistence", ["SQLite based historical tracking"]),
    ("Alerting", ["Conditional email notifications"])
])

# Slide 6: Dashboard Feature
add_slide(prs, 1, "Feature: Rich Dashboard", [
    "A comprehensive Streamlit UI providing 8 distinct sub-pages:",
    ("1. Home", ["Entry point for configuring target and API keys"]),
    ("2. Summary", ["High-level overview of the scan result"]),
    ("3. Analysis", ["Deep dive into findings and VirusTotal flags"]),
    ("4. Visuals", ["Plotly-powered interactive charts"]),
    ("5. History", ["Tracking past scans and trends"]),
    ("6. Recommendations", ["Actionable mitigation steps"]),
    ("7. Risk Map", ["Visual breakdown of the total risk"]),
    ("8. System Info", ["Target metadata details"])
])

# Slide 7: Nmap Discovery Feature
add_slide(prs, 1, "Feature: Open-Port Discovery", [
    "Powered by python-nmap on the backend.",
    "Actively scans a given domain or IP address for exposed services.",
    "Extracts port numbers, protocols, states, and product versions.",
    "Identifies exactly where the attack surface begins."
])

# Slide 8: VT Enrichment Feature
add_slide(prs, 1, "Feature: Threat Intelligence", [
    "Deep integration with VirusTotal API.",
    "Fetches global reputation data for Domains, IPs, URLs, and File Hashes.",
    "Identifies if the asset is flagged as malicious, suspicious, or clean by multiple vendors.",
    "Retrieves contextual community votes and historical analysis scores."
])

# Slide 9: Risk Scoring Feature
add_slide(prs, 1, "Feature: Structured Risk Scoring", [
    "Mathematical model to reduce subjective guessing.",
    "Exposure Score: Based on open ports and services.",
    "Threat Score: Derived directly from VT malicious flags.",
    "Context Score: Adjustments based on target type and history.",
    "Overall Risk: A consolidated score between 0 and 100 representing urgency."
])

# Slide 10: Persistence Feature
add_slide(prs, 1, "Feature: Data Persistence", [
    "Built-in robust SQLite database.",
    "Tracks every scan along with its specific timestamp.",
    "Saves all 4 risk dimensions (Exposure, Threat, Context, Risk).",
    "Allows analysts to see if a system's security posture is improving or degrading over time."
])

# Slide 11: Alerting Feature
add_slide(prs, 1, "Feature: Automated Email Alerts", [
    "Asynchronous alerting system integrated into the backend.",
    "Triggers notifications immediately if an asset contains malicious detections.",
    "Also triggers if the Overall Risk score exceeds the critical threshold (70).",
    "Delivers a summary right to the analyst's inbox for immediate response."
])

# Slide 12: Supported Targets
add_slide(prs, 1, "Supported Targets", [
    "Designed to run comprehensive workflows against:",
    "1. Domains (e.g., example.com)",
    "2. IP Addresses (e.g., 8.8.8.8)",
    "3. URLs",
    "4. File Hashes (SHA-256 for historical intelligence)"
])

# Slide 13: Architecture
add_slide(prs, 1, "Architecture Workflow", [
    "1. Analyst inputs target, API keys, and email in Streamlit UI.",
    "2. UI calls FastAPI backend `/scan/{target}`.",
    "3. Backend delegates discovery to Nmap.",
    "4. Backend delegates intel gathering to VirusTotal API.",
    "5. Normalizer aggregates and Scoring Engine calculates risk.",
    "6. Results saved to SQLite DB.",
    "7. Alert Evaluator dispatches email if conditions met.",
    "8. UI visualizes final structured JSON payload."
])

# Slide 14: Future Enhancements
add_slide(prs, 1, "Future Enhancements", [
    "WHOIS Enrichment: Gathering domain ownership context.",
    "SSL/TLS Validation: Checking for expired or misconfigured certificates.",
    "File Upload Analysis: Detonating raw binaries/files directly.",
    "Export Engine: Generating automated PDF/HTML executive reports."
])

# Slide 15: Conclusion
add_slide(prs, 0, "Thank You", [
    "Questions?"
])

# Save the presentation
output_pptx = "CyberScan_Presentation.pptx"
prs.save(output_pptx)
print(f"✅ Real PPTX Presentation generated successfully at: {output_pptx}")
